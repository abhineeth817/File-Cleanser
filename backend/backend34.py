import os # For file path operations
import io# For in-memory byte streams
import spacy#   for NLP tasks
import fitz  # PyMuPDF
import docx#  For DOCX files
import openpyxl#  For XLSX files
import pptx#  For PPTX files
import pytesseract#  For OCR
from PIL import Image, ImageDraw, ImageFont#  For image processing
import textwrap#  For text wrapping

from pptx.enum.shapes import MSO_SHAPE_TYPE#  For identifying shape types in PPTX
from docx import Document#  For creating DOCX files
from presidio_analyzer import AnalyzerEngine#   for PII detection
from presidio_anonymizer import AnonymizerEngine#   for PII anonymization
from reportlab.pdfgen import canvas#  For PDF generation
from reportlab.lib.pagesizes import letter#  For PDF page sizes

#Presidio Analyzer and Anonymizer
analyzer = AnalyzerEngine()
anonymizer = AnonymizerEngine()

#spaCy model
try:
    nlp = spacy.load("en_core_web_sm")
except OSError:
    print("Downloading 'en_core_web_sm' model...")
    spacy.cli.download("en_core_web_sm")
    nlp = spacy.load("en_core_web_sm")
#for cleansing
def cleanse_text(text):
    if not text or text.isspace():
        return ""
    results = analyzer.analyze(text=text, language="en")
    anonymized_result = anonymizer.anonymize(text=text, analyzer_results=results)
    return anonymized_result.text

#extracting pdf
def extract_pdf(file_path):
    doc = fitz.open(file_path)
    full_text = []
    for page_num in range(len(doc)):
        page_content = [f"\n\n--- PAGE {page_num + 1} ---\n"]
        page = doc.load_page(page_num)
        
        page_text = page.get_text("text")
        if page_text.strip():
            page_content.append(page_text)
        
        image_list = page.get_images(full=True)
        if image_list:
            ocr_texts = []
            for img_index, img in enumerate(image_list):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image = Image.open(io.BytesIO(image_bytes))
                    ocr_text = pytesseract.image_to_string(image)
                    if ocr_text.strip():
                        ocr_texts.append(f"\n--- OCR from Image {img_index + 1} ---\n{ocr_text}\n")
                except Exception as e:
                    print(f"Warning: Could not process image {img_index + 1} on PDF page {page_num+1}. Error: {e}")
            
            if ocr_texts:
                page_content.append("\n--- OCR TEXT FROM IMAGES ON THIS PAGE ---")
                page_content.extend(ocr_texts)

        full_text.append("".join(page_content))

    return "\n".join(full_text)

#extracting docx or word files
def extract_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

#extracting xlsx or excel files
def extract_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path)
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
    return text

#extracting pptx or powerpoint files
def extract_pptx(file_path):
    prs = pptx.Presentation(file_path)
    full_text = []
    for i, slide in enumerate(prs.slides):
        full_text.append(f"\n\n--- SLIDE {i + 1} ---\n")
        slide_texts = []
        slide_ocr_texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
            
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE and shape.has_table:
                table_text = []
                tbl = shape.table
                for row in tbl.rows:
                    row_text = [cell.text for cell in row.cells]
                    table_text.append("\t".join(row_text))
                if any(table_text):
                    slide_texts.append("\n".join(table_text))
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_bytes = shape.image.blob
                    image = Image.open(io.BytesIO(image_bytes))
                    ocr_text = pytesseract.image_to_string(image)
                    if ocr_text.strip():
                       slide_ocr_texts.append(f"--- OCR from Image on Slide {i+1} ---\n{ocr_text}")
                except Exception as e:
                    print(f"Warning: Could not process image on slide {i+1}. Error: {e}")

        if slide_texts:
            full_text.append("\n".join(slide_texts))
        if slide_ocr_texts:
            full_text.append("\n--- OCR TEXT FROM IMAGES ON THIS SLIDE ---\n")
            full_text.append("\n\n".join(slide_ocr_texts))

    return "\n".join(full_text)

#extracting text from image files
def extract_image(file_path):
    try:
        image = Image.open(file_path)
        return pytesseract.image_to_string(image)
    except Exception as e:
        print(f"Error while processing image file {file_path}: {e}")
        return ""
    
#saving pdf
def save_pdf(cleansed, output_path):
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    margin = 50
    
    textobject = c.beginText()
    textobject.setTextOrigin(margin, height - margin)
    textobject.setFont("Helvetica", 10)

    for line in cleansed.split('\n'):
        wrapped_lines = textwrap.wrap(line, width=100)
        if not wrapped_lines:
            textobject.textLine('')
        
        for wrapped_line in wrapped_lines:
            textobject.textLine(wrapped_line)
        
        if textobject.getY() < margin:
            c.drawText(textobject)
            c.showPage()
            textobject = c.beginText()
            textobject.setTextOrigin(margin, height - margin)
            textobject.setFont("Helvetica", 10)
            
    c.drawText(textobject)
    c.save()

#saving word files
def save_docx(cleansed, output_path):
    doc = Document()
    for line in cleansed.split("\n"):
        doc.add_paragraph(line)
    doc.save(output_path)

#saving excel files
def save_xlsx(cleansed, output_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    for line in cleansed.split("\n"):
        ws.append([line])
    wb.save(output_path)

#saving ppt files
def save_pptx(cleansed, output_path):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[5]  
    slide = prs.slides.add_slide(slide_layout)
    
    left = top = width = height = pptx.util.Inches(1)
    txBox = slide.shapes.add_textbox(left, top, prs.slide_width - (2*width), prs.slide_height - (2*height))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = cleansed
    
    prs.save(output_path)

#saving image
def save_image(cleansed, output_path):
    width, height = 1200, 1600  
    background_color = "white"
    text_color = "black"
    margin = 50
    
    try:
        font = ImageFont.truetype("arial.ttf", 24)
    except IOError:
        font = ImageFont.load_default()

    image = Image.new("RGB", (width, height), background_color)
    draw = ImageDraw.Draw(image)

    wrapped_text = "\n".join(textwrap.wrap(cleansed, width=100, replace_whitespace=False))
    
    draw.multiline_text((margin, margin), wrapped_text, font=font, fill=text_color)
    image.save(output_path)

#saving text files
def save_txt(cleansed, output_path):
    with open(output_path, "w", encoding='utf-8') as f:
        f.write(cleansed)

#dispatcher
def process_file(file_path, output_dir="output", preserve_file_type=True):
    if not os.path.exists(file_path):
        print(f"Error: File not found at {file_path}")
        return None, None

    os.makedirs(output_dir, exist_ok=True)
    ext = os.path.splitext(file_path)[-1].lower()
    raw_text = ""
    try:
        if ext == ".pdf":
            raw_text = extract_pdf(file_path)
        elif ext == ".docx":
            raw_text = extract_docx(file_path)
        elif ext == ".xlsx":
            raw_text = extract_xlsx(file_path)
        elif ext == ".pptx":
            raw_text = extract_pptx(file_path)
        elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
            raw_text = extract_image(file_path)
        else:
            print(f"Unsupported file format: {ext}")
            return None, None
    except Exception as e:
        print(f"An error occurred during text extraction from {file_path}: {e}")
        return None, None

    cleansed_text = cleanse_text(raw_text)
    
    filename = os.path.basename(file_path)
    output_path = None
    try:
        if preserve_file_type:
            if ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
                output_filename = f"cleansed_{os.path.splitext(filename)[0]}.png"
                output_path = os.path.join(output_dir, output_filename)
                save_image(cleansed_text, output_path)
            else:
                output_filename = f"cleansed_{filename}"
                output_path = os.path.join(output_dir, output_filename)
                if ext == ".pdf":
                    save_pdf(cleansed_text, output_path)
                elif ext == ".docx":
                    save_docx(cleansed_text, output_path)
                elif ext == ".xlsx":
                    save_xlsx(cleansed_text, output_path)
                elif ext == ".pptx":
                    save_pptx(cleansed_text, output_path)
        else:
            output_filename = f"cleansed_{os.path.splitext(filename)[0]}.txt"
            output_path = os.path.join(output_dir, output_filename)
            save_txt(cleansed_text, output_path)
    except Exception as e:
        print(f"An error occurred while saving the cleansed file: {e}")
        return cleansed_text, None 

    return cleansed_text, output_path

#main
if __name__ == "__main__":
    file = "/users/work/documents/File_012.pdf"  
    preserve = True
    
    print(f"Processing file: {file}...")
    print(f"Preserve file type: {preserve}")

    cleansed_content, saved_file = process_file(
        file, 
        preserve_file_type=preserve
    )

    if cleansed_content and saved_file:
        print(f"\nSuccessfully processed file.")
        print(f"Cleansed file saved at: {saved_file}")
    else:
        print("\nFile processing failed.")

